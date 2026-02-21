"""PPTX extractor implementation."""

from __future__ import annotations

from pathlib import Path
from uuid import uuid4

from pptx import Presentation

from app.services.extraction.base import BaseExtractor
from app.services.extraction.models import ExtractionWarning, Segment


class PptxExtractor(BaseExtractor):
    def extract(
        self,
        file_path: Path,
        source_path: str,
        doc_uid: str,
        extracted_at: str,
    ) -> tuple[list[Segment], list[ExtractionWarning]]:
        presentation = Presentation(str(file_path))
        segments: list[Segment] = []
        warnings: list[ExtractionWarning] = []

        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_fragments: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    text = shape.text_frame.text or ""
                    if text:
                        text_fragments.append(text)

            merged_text = "\n".join(text_fragments)
            locator = {"slide": slide_index}
            segments.append(
                Segment(
                    segment_id=str(uuid4()),
                    doc_uid=doc_uid,
                    source_path=source_path,
                    file_type="pptx",
                    unit_type="slide",
                    unit_index=slide_index,
                    text=merged_text,
                    locator=locator,
                    extracted_at=extracted_at,
                )
            )
            if not merged_text.strip():
                warnings.append(
                    ExtractionWarning(
                        code="empty_text_slide",
                        message="Slide contains no extractable text.",
                        locator=locator,
                    )
                )

        if not segments:
            warnings.append(
                ExtractionWarning(
                    code="pptx_no_slides",
                    message="PPTX contains no slides.",
                )
            )

        return segments, warnings
