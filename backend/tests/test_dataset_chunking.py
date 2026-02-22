# backend/tests/test_dataset_chunking.py
import hashlib
import json
import subprocess
import sys
from pathlib import Path

import yaml
from docx import Document
from pptx import Presentation
from pptx.util import Inches


BASE = Path(__file__).resolve().parents[1]
SCRIPT = BASE / "scripts" / "datasets" / "build_chunks.py"


def _run_build_chunks(command: list[str]) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, cwd=BASE, capture_output=True, text=True, check=False)


def _build_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("Lecture 01", style="Heading 1")
    for index in range(6):
        document.add_paragraph(
            f"Docx paragraph {index}. This is sample content for chunk testing with traceability."
        )
    document.save(path)


def _build_pptx(path: Path) -> None:
    presentation = Presentation()
    blank_layout = (
        presentation.slide_layouts[6]
        if len(presentation.slide_layouts) > 6
        else presentation.slide_layouts[-1]
    )
    for index in range(2):
        slide = presentation.slides.add_slide(blank_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        textbox.text_frame.text = (
            f"Slide {index + 1} text. Chunk traceability should preserve slide boundaries."
        )
    presentation.save(path)


def _write_manifest(manifest_path: Path, datasets: list[dict]) -> None:
    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(
        yaml.safe_dump({"version": "v1", "datasets": datasets}, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )


def _read_jsonl(path: Path) -> list[dict]:
    lines = [line for line in path.read_text(encoding="utf-8").splitlines() if line.strip()]
    return [json.loads(line) for line in lines]


def _sha256_file(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def test_build_chunks_outputs_txt_docx_pptx_and_stats(tmp_path) -> None:
    backend_root = tmp_path / "backend"
    raw_dir = backend_root / "datasets_local" / "raw"
    redacted_dir = backend_root / "datasets_local" / "redacted"
    out_dir = backend_root / "datasets_local" / "exports" / "chunks"
    raw_dir.mkdir(parents=True, exist_ok=True)
    redacted_dir.mkdir(parents=True, exist_ok=True)

    txt_raw_path = raw_dir / "notes.txt"
    txt_raw_path.write_text("RAW SHOULD NOT BE USED", encoding="utf-8")
    txt_redacted_path = redacted_dir / "notes.txt"
    txt_redacted_path.write_text(
        (
            "REDACTED paragraph A. This paragraph is intentionally long for chunking behavior.\n\n"
            "REDACTED paragraph B. This paragraph also has enough text to trigger merge behavior.\n\n"
            "REDACTED paragraph C. Keep this paragraph for verifying max-chars constraints.\n\n"
            "REDACTED paragraph D. Final paragraph for sequence and stats validation."
        ),
        encoding="utf-8",
    )

    docx_path = raw_dir / "lecture.docx"
    pptx_path = raw_dir / "slides.pptx"
    _build_docx(docx_path)
    _build_pptx(pptx_path)

    manifest_path = backend_root / "docs" / "ai" / "datasets" / "manifest.v1.yaml"
    _write_manifest(
        manifest_path,
        datasets=[
            {
                "dataset_id": "ds-txt",
                "allowed_use": "infer_only",
                "license": {"type": "TBD", "evidence": ""},
                "privacy": {"redaction_status": "pending"},
                "files": [
                    {
                        "relative_path": "backend/datasets_local/raw/notes.txt",
                        "file_type": "text",
                        "sha256": "0" * 64,
                        "size_bytes": txt_raw_path.stat().st_size,
                    }
                ],
                "updated_at": "2026-02-22T00:00:00Z",
            },
            {
                "dataset_id": "ds-docx",
                "allowed_use": "infer_only",
                "license": {"type": "TBD", "evidence": ""},
                "privacy": {"redaction_status": "pending"},
                "files": [
                    {
                        "relative_path": "backend/datasets_local/raw/lecture.docx",
                        "file_type": "docx",
                        "sha256": "0" * 64,
                        "size_bytes": docx_path.stat().st_size,
                    }
                ],
                "updated_at": "2026-02-22T00:00:00Z",
            },
            {
                "dataset_id": "ds-pptx",
                "allowed_use": "infer_only",
                "license": {"type": "TBD", "evidence": ""},
                "privacy": {"redaction_status": "pending"},
                "files": [
                    {
                        "relative_path": "backend/datasets_local/raw/slides.pptx",
                        "file_type": "pptx",
                        "sha256": "0" * 64,
                        "size_bytes": pptx_path.stat().st_size,
                    }
                ],
                "updated_at": "2026-02-22T00:00:00Z",
            },
        ],
    )

    result = _run_build_chunks(
        [
            sys.executable,
            str(SCRIPT),
            "--manifest",
            str(manifest_path),
            "--out-dir",
            str(out_dir),
            "--max-chars",
            "120",
            "--min-chars",
            "40",
        ]
    )
    assert result.returncode == 0, result.stderr

    txt_records = _read_jsonl(out_dir / "ds-txt.chunks.v1.jsonl")
    assert txt_records
    assert [record["chunk_id"] for record in txt_records] == [
        f"ch-ds-txt-{index:06d}" for index in range(1, len(txt_records) + 1)
    ]
    assert all(len(record["text"]) <= 120 for record in txt_records)
    assert all("paragraph_start" in record["meta"] for record in txt_records)
    assert all("paragraph_end" in record["meta"] for record in txt_records)
    assert all(record["meta"]["source_relative_path"] == "backend/datasets_local/raw/notes.txt" for record in txt_records)
    assert all(record["meta"]["sha256"] == _sha256_file(txt_redacted_path) for record in txt_records)
    assert any("REDACTED paragraph" in record["text"] for record in txt_records)
    assert all("RAW SHOULD NOT BE USED" not in record["text"] for record in txt_records)

    docx_records = _read_jsonl(out_dir / "ds-docx.chunks.v1.jsonl")
    assert docx_records
    assert all("paragraph_start" in record["meta"] for record in docx_records)
    assert all("paragraph_end" in record["meta"] for record in docx_records)
    assert min(record["meta"]["paragraph_start"] for record in docx_records) == 0
    assert max(record["meta"]["paragraph_end"] for record in docx_records) >= 5

    pptx_records = _read_jsonl(out_dir / "ds-pptx.chunks.v1.jsonl")
    assert pptx_records
    assert all("slide_start" in record["meta"] for record in pptx_records)
    assert all("slide_end" in record["meta"] for record in pptx_records)
    assert min(record["meta"]["slide_start"] for record in pptx_records) == 1
    assert max(record["meta"]["slide_end"] for record in pptx_records) >= 2

    stats_path = out_dir / "chunk_stats.v1.json"
    stats = json.loads(stats_path.read_text(encoding="utf-8"))
    assert "datasets" in stats
    assert stats["datasets"]["ds-txt"]["chunk_count"] == len(txt_records)
    assert stats["datasets"]["ds-docx"]["chunk_count"] == len(docx_records)
    assert stats["datasets"]["ds-pptx"]["chunk_count"] == len(pptx_records)
    assert stats["datasets"]["ds-txt"]["total_chars"] >= stats["datasets"]["ds-txt"]["chunk_count"]


def test_build_chunks_dry_run_does_not_write_files(tmp_path) -> None:
    backend_root = tmp_path / "backend"
    raw_dir = backend_root / "datasets_local" / "raw"
    raw_dir.mkdir(parents=True, exist_ok=True)

    txt_path = raw_dir / "dry-run.txt"
    txt_path.write_text("Dry run paragraph one.\n\nDry run paragraph two.", encoding="utf-8")

    manifest_path = backend_root / "docs" / "ai" / "datasets" / "manifest.v1.yaml"
    _write_manifest(
        manifest_path,
        datasets=[
            {
                "dataset_id": "ds-dry-run",
                "allowed_use": "infer_only",
                "license": {"type": "TBD", "evidence": ""},
                "privacy": {"redaction_status": "pending"},
                "files": [
                    {
                        "relative_path": "backend/datasets_local/raw/dry-run.txt",
                        "file_type": "text",
                        "sha256": "0" * 64,
                        "size_bytes": txt_path.stat().st_size,
                    }
                ],
                "updated_at": "2026-02-22T00:00:00Z",
            }
        ],
    )

    out_dir = backend_root / "datasets_local" / "exports" / "chunks"
    result = _run_build_chunks(
        [
            sys.executable,
            str(SCRIPT),
            "--manifest",
            str(manifest_path),
            "--out-dir",
            str(out_dir),
            "--dry-run",
        ]
    )
    assert result.returncode == 0, result.stderr
    assert not out_dir.exists()
