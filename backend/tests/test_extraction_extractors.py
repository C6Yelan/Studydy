import json
from pathlib import Path

from docx import Document
from pypdf import PdfWriter
from pptx import Presentation
from pptx.util import Inches

from app.services.extraction.pipeline import extract_document


def _build_pdf(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    with path.open("wb") as file_obj:
        writer.write(file_obj)


def _build_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("DOCX paragraph one")
    document.add_paragraph("DOCX paragraph two")
    document.save(path)


def _build_pptx(path: Path) -> None:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]

    slide_1 = presentation.slides.add_slide(blank_layout)
    textbox_1 = slide_1.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox_1.text_frame.text = "PPTX slide one text"

    slide_2 = presentation.slides.add_slide(blank_layout)
    textbox_2 = slide_2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox_2.text_frame.text = "PPTX slide two text"

    presentation.save(path)


def test_pdf_extractor_returns_page_segments_and_warnings(tmp_path) -> None:
    pdf_path = tmp_path / "sample.pdf"
    _build_pdf(pdf_path)

    result = extract_document(pdf_path, run_id="test-run")

    assert result.status == "success"
    assert result.file_type == "pdf"
    assert result.doc_uid == result.sha256
    assert len(result.segments) == 1
    assert result.segments[0].unit_type == "page"
    assert result.segments[0].unit_index == 1
    assert result.segments[0].locator == {"page": 1}
    assert result.segments[0].text == ""
    assert any(warning.code == "empty_text_page" for warning in result.warnings)

    rerun = extract_document(pdf_path, run_id="test-run-rerun")
    assert rerun.doc_uid == result.doc_uid


def test_docx_extractor_returns_paragraph_segments(tmp_path) -> None:
    docx_path = tmp_path / "sample.docx"
    _build_docx(docx_path)

    result = extract_document(docx_path, run_id="test-run")

    assert result.status == "success"
    assert result.file_type == "docx"
    assert result.doc_uid == result.sha256
    assert [segment.unit_index for segment in result.segments] == [1, 2]
    assert [segment.locator for segment in result.segments] == [{"paragraph": 1}, {"paragraph": 2}]
    assert [segment.text for segment in result.segments] == ["DOCX paragraph one", "DOCX paragraph two"]
    assert result.warnings == []


def test_pptx_extractor_returns_slide_segments(tmp_path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _build_pptx(pptx_path)

    result = extract_document(pptx_path, run_id="test-run")

    assert result.status == "success"
    assert result.file_type == "pptx"
    assert result.doc_uid == result.sha256
    assert [segment.unit_index for segment in result.segments] == [1, 2]
    assert [segment.locator for segment in result.segments] == [{"slide": 1}, {"slide": 2}]
    assert "PPTX slide one text" in result.segments[0].text
    assert "PPTX slide two text" in result.segments[1].text
    assert result.warnings == []



def test_extract_result_serialization_is_json_safe(tmp_path) -> None:
    docx_path = tmp_path / "sample.docx"
    _build_docx(docx_path)

    result = extract_document(docx_path, run_id="test-run")
    payload = result.to_report_dict()

    serialized = json.dumps(payload)
    assert serialized
