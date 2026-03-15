# backend/scripts/datasets/build_chunks.py
"""Build dataset chunks (T3) from manifest entries into intermediate JSONL."""

from __future__ import annotations

import argparse
import hashlib
import json
import sys
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import yaml
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


BACKEND_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_MANIFEST_PATH = BACKEND_ROOT / "docs" / "ai" / "datasets" / "manifest.v1.yaml"
DEFAULT_OUT_DIR = BACKEND_ROOT / "datasets_local" / "exports" / "chunks"

SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".txt", ".md"}
PUNCTUATION_BREAKS = {"。", "！", "？", ".", "!", "?"}


@dataclass(slots=True, frozen=True)
class Unit:
    text: str
    start: int
    end: int
    heading: str | None = None
    title: str | None = None


@dataclass(slots=True, frozen=True)
class ChunkSpan:
    text: str
    start: int
    end: int
    heading: str | None = None
    title: str | None = None


@dataclass(slots=True, frozen=True)
class ExtractedFile:
    units: list[Unit]
    locator_key: str
    heading: str | None
    title: str | None


def utc_now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def load_manifest(path: Path) -> dict[str, Any]:
    data = yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    if not isinstance(data, dict):
        raise ValueError(f"Manifest must be a mapping object: {path}")
    return data


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as file_obj:
        while True:
            chunk = file_obj.read(1024 * 1024)
            if not chunk:
                break
            digest.update(chunk)
    return digest.hexdigest()


def _build_document_id(*, dataset_id: str, source_relative_path: str, sha256: str) -> str:
    digest = hashlib.sha256(
        f"{dataset_id}:{source_relative_path}:{sha256}".encode("utf-8")
    ).hexdigest()[:16]
    return f"doc-{digest}"


def _infer_backend_root(manifest_path: Path) -> Path:
    for candidate in [manifest_path.parent, *manifest_path.parents]:
        if candidate.name == "backend":
            return candidate.resolve()
    return BACKEND_ROOT


def _candidate_roots(manifest_path: Path) -> list[Path]:
    backend_root = _infer_backend_root(manifest_path)
    roots = [
        manifest_path.parent.resolve(),
        backend_root,
        backend_root.parent.resolve(),
        BACKEND_ROOT,
        BACKEND_ROOT.parent.resolve(),
        Path.cwd().resolve(),
    ]

    unique: list[Path] = []
    seen: set[str] = set()
    for root in roots:
        key = root.as_posix()
        if key in seen:
            continue
        seen.add(key)
        unique.append(root)
    return unique


def _resolve_manifest_entry_path(manifest_path: Path, path_text: str) -> Path:
    path_obj = Path(path_text)
    if path_obj.is_absolute():
        return path_obj.resolve()

    for root in _candidate_roots(manifest_path):
        candidate = (root / path_obj).resolve()
        if candidate.exists():
            return candidate

    return (_candidate_roots(manifest_path)[0] / path_obj).resolve()


def _find_redacted_file(raw_path: Path, backend_root: Path) -> Path:
    redacted_root = backend_root / "datasets_local" / "redacted"
    raw_root = backend_root / "datasets_local" / "raw"
    if not redacted_root.exists():
        return raw_path

    try:
        relative_under_raw = raw_path.resolve().relative_to(raw_root.resolve())
        candidate = (redacted_root / relative_under_raw).resolve()
        if candidate.exists() and candidate.is_file():
            return candidate
    except ValueError:
        pass

    direct_candidate = (redacted_root / raw_path.name).resolve()
    if direct_candidate.exists() and direct_candidate.is_file():
        return direct_candidate

    same_name = sorted(path for path in redacted_root.rglob(raw_path.name) if path.is_file())
    if same_name:
        return same_name[0].resolve()

    return raw_path.resolve()


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="latin-1")


def _extract_text_units(path: Path, *, markdown: bool) -> ExtractedFile:
    content = _read_text_file(path)
    lines = content.splitlines()

    title: str | None = None
    heading: str | None = None
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if title is None:
            title = stripped.lstrip("# ").strip() if stripped.startswith("#") else stripped
        if markdown and heading is None and stripped.startswith("#"):
            heading = stripped.lstrip("# ").strip()
        if title is not None and (heading is not None or not markdown):
            break

    units: list[Unit] = []
    paragraph_lines: list[str] = []
    paragraph_index = 0

    def flush_paragraph() -> None:
        nonlocal paragraph_index
        if not paragraph_lines:
            return
        text = "\n".join(paragraph_lines).strip()
        units.append(Unit(text=text, start=paragraph_index, end=paragraph_index, heading=heading, title=title))
        paragraph_lines.clear()
        paragraph_index += 1

    for line in lines:
        if not line.strip():
            flush_paragraph()
            continue
        paragraph_lines.append(line.rstrip())

    flush_paragraph()
    return ExtractedFile(units=units, locator_key="paragraph", heading=heading, title=title)


def _extract_docx_units(path: Path) -> ExtractedFile:
    document = Document(str(path))
    heading: str | None = None
    title: str | None = None
    units: list[Unit] = []

    for paragraph_index, paragraph in enumerate(document.paragraphs):
        text = paragraph.text or ""
        stripped = text.strip()
        style_name = (paragraph.style.name or "").strip().lower() if paragraph.style is not None else ""

        if title is None and stripped:
            title = stripped
        if heading is None and stripped and style_name.startswith("heading"):
            heading = stripped

        units.append(
            Unit(
                text=text,
                start=paragraph_index,
                end=paragraph_index,
                heading=heading,
                title=title,
            )
        )

    return ExtractedFile(units=units, locator_key="paragraph", heading=heading, title=title)


def _extract_pptx_units(path: Path) -> ExtractedFile:
    presentation = Presentation(str(path))
    heading: str | None = None
    title = (presentation.core_properties.title or "").strip() or None
    units: list[Unit] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        fragments: list[str] = []
        slide_heading: str | None = None

        if slide.shapes.title is not None:
            candidate = (slide.shapes.title.text or "").strip()
            if candidate:
                slide_heading = candidate
                if heading is None:
                    heading = candidate
                if title is None:
                    title = candidate

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue
            text = shape.text_frame.text or ""
            if text.strip():
                fragments.append(text.strip())

        units.append(
            Unit(
                text="\n".join(fragments),
                start=slide_index,
                end=slide_index,
                heading=slide_heading or heading,
                title=title,
            )
        )

    return ExtractedFile(units=units, locator_key="slide", heading=heading, title=title)


def _extract_pdf_units(path: Path) -> ExtractedFile:
    reader = PdfReader(str(path))
    metadata_title = None
    if reader.metadata is not None:
        metadata_title = str(reader.metadata.title or "").strip() or None

    units: list[Unit] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        units.append(Unit(text=text, start=page_index, end=page_index, title=metadata_title))

    return ExtractedFile(units=units, locator_key="page", heading=None, title=metadata_title)


def extract_units(path: Path) -> ExtractedFile:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_units(path)
    if suffix == ".docx":
        return _extract_docx_units(path)
    if suffix == ".pptx":
        return _extract_pptx_units(path)
    if suffix == ".txt":
        return _extract_text_units(path, markdown=False)
    if suffix == ".md":
        return _extract_text_units(path, markdown=True)
    raise ValueError(f"Unsupported file suffix for chunking: {path.suffix}")


def _find_cut_index(window: str, min_chars: int) -> int:
    lower_bound = max(1, min_chars)
    for index in range(len(window) - 1, lower_bound - 1, -1):
        if window[index] == "\n":
            return index + 1
    for index in range(len(window) - 1, lower_bound - 1, -1):
        if window[index] in PUNCTUATION_BREAKS:
            return index + 1
    return len(window)


def split_long_text(text: str, max_chars: int, min_chars: int) -> list[str]:
    cleaned = text.strip()
    if not cleaned:
        return []
    if len(cleaned) <= max_chars:
        return [cleaned]

    result: list[str] = []
    cursor = 0
    while cursor < len(cleaned):
        remaining = cleaned[cursor:]
        if len(remaining) <= max_chars:
            piece = remaining.strip()
            if piece:
                result.append(piece)
            break

        window = cleaned[cursor : cursor + max_chars]
        cut = _find_cut_index(window, min_chars=min_chars)
        piece = cleaned[cursor : cursor + cut].strip()
        if not piece:
            cut = max_chars
            piece = cleaned[cursor : cursor + cut].strip()
        if piece:
            result.append(piece)
        cursor += cut

    if len(result) >= 2 and len(result[-1]) < min_chars:
        merged = f"{result[-2]}\n{result[-1]}".strip()
        if len(merged) <= max_chars:
            result[-2] = merged
            result.pop()

    return result


def build_spans(
    units: list[Unit],
    *,
    max_chars: int,
    min_chars: int,
    overlap_chars: int,
) -> list[ChunkSpan]:
    if max_chars <= 0:
        raise ValueError("max_chars must be > 0")
    if min_chars < 0:
        raise ValueError("min_chars must be >= 0")
    if min_chars > max_chars:
        raise ValueError("min_chars cannot be greater than max_chars")
    if overlap_chars < 0:
        raise ValueError("overlap_chars must be >= 0")

    segments: list[ChunkSpan] = []
    for unit in units:
        unit_text = (unit.text or "").strip()
        if not unit_text:
            continue
        pieces = split_long_text(unit_text, max_chars=max_chars, min_chars=min_chars)
        for piece in pieces:
            segments.append(
                ChunkSpan(
                    text=piece,
                    start=unit.start,
                    end=unit.end,
                    heading=unit.heading,
                    title=unit.title,
                )
            )

    chunks: list[ChunkSpan] = []
    current: ChunkSpan | None = None

    for segment in segments:
        if current is None:
            current = segment
            continue

        merged_text = f"{current.text}\n\n{segment.text}"
        if len(merged_text) <= max_chars:
            current = ChunkSpan(
                text=merged_text,
                start=current.start,
                end=segment.end,
                heading=current.heading or segment.heading,
                title=current.title or segment.title,
            )
        else:
            chunks.append(current)
            current = segment

    if current is not None:
        chunks.append(current)

    if min_chars > 0 and len(chunks) >= 2:
        index = 0
        while index < len(chunks) - 1:
            if len(chunks[index].text) < min_chars:
                merged_text = f"{chunks[index].text}\n\n{chunks[index + 1].text}"
                if len(merged_text) <= max_chars:
                    chunks[index] = ChunkSpan(
                        text=merged_text,
                        start=chunks[index].start,
                        end=chunks[index + 1].end,
                        heading=chunks[index].heading or chunks[index + 1].heading,
                        title=chunks[index].title or chunks[index + 1].title,
                    )
                    chunks.pop(index + 1)
                    continue
            index += 1

        if len(chunks) >= 2 and len(chunks[-1].text) < min_chars:
            merged_tail = f"{chunks[-2].text}\n\n{chunks[-1].text}"
            if len(merged_tail) <= max_chars:
                chunks[-2] = ChunkSpan(
                    text=merged_tail,
                    start=chunks[-2].start,
                    end=chunks[-1].end,
                    heading=chunks[-2].heading or chunks[-1].heading,
                    title=chunks[-2].title or chunks[-1].title,
                )
                chunks.pop()

    if overlap_chars > 0 and len(chunks) >= 2:
        overlapped: list[ChunkSpan] = [chunks[0]]
        for index in range(1, len(chunks)):
            previous = overlapped[-1]
            current_chunk = chunks[index]
            overlap = previous.text[-overlap_chars:] if previous.text else ""
            if overlap:
                candidate = f"{overlap}\n{current_chunk.text}"
                text = candidate if len(candidate) <= max_chars else current_chunk.text
            else:
                text = current_chunk.text
            overlapped.append(
                ChunkSpan(
                    text=text,
                    start=current_chunk.start,
                    end=current_chunk.end,
                    heading=current_chunk.heading,
                    title=current_chunk.title,
                )
            )
        chunks = overlapped

    return chunks


def _ensure_supported(path: Path) -> None:
    if path.suffix.lower() not in SUPPORTED_SUFFIXES:
        raise ValueError(f"Unsupported file type for T3 chunking: {path}")


def _build_locator(
    *,
    locator_key: str,
    span: ChunkSpan,
    source_relative_path: str,
    document_id: str,
    chunk_start_char: int,
    chunk_end_char: int,
) -> dict[str, Any]:
    locator: dict[str, Any] = {
        "source_file": source_relative_path,
        "document_id": document_id,
        "chunk_start_char": chunk_start_char,
        "chunk_end_char": chunk_end_char,
    }

    if locator_key == "page":
        locator["page_start"] = span.start
        locator["page_end"] = span.end
        if span.start == span.end:
            locator["page"] = span.start
    elif locator_key == "paragraph":
        locator["paragraph_start"] = span.start
        locator["paragraph_end"] = span.end
        if span.start == span.end:
            locator["paragraph_index"] = span.start
    elif locator_key == "slide":
        locator["slide_index"] = span.start
        if span.start != span.end:
            locator["slide_start"] = span.start
            locator["slide_end"] = span.end

    return locator


def build_dataset_chunks(
    *,
    manifest_path: Path,
    max_chars: int,
    min_chars: int,
    overlap_chars: int,
    dataset_id_filter: str | None,
) -> tuple[dict[str, list[dict[str, Any]]], dict[str, Any]]:
    manifest = load_manifest(manifest_path)
    datasets = manifest.get("datasets", [])
    if not isinstance(datasets, list):
        raise ValueError("Manifest field `datasets` must be an array")

    selected = []
    for dataset in datasets:
        if not isinstance(dataset, dict):
            continue
        dataset_id = dataset.get("dataset_id")
        if not isinstance(dataset_id, str) or not dataset_id.strip():
            continue
        if dataset_id_filter is not None and dataset_id != dataset_id_filter:
            continue
        selected.append(dataset)

    if dataset_id_filter is not None and not selected:
        raise ValueError(f"Dataset id not found in manifest: {dataset_id_filter}")

    backend_root = _infer_backend_root(manifest_path)
    created_at = utc_now_iso()
    dataset_chunks: dict[str, list[dict[str, Any]]] = {}
    per_dataset_stats: dict[str, dict[str, float | int]] = {}

    for dataset in selected:
        dataset_id = str(dataset["dataset_id"])
        files = dataset.get("files", [])
        if not isinstance(files, list):
            files = []

        records: list[dict[str, Any]] = []
        sequence = 1

        for file_entry in files:
            if not isinstance(file_entry, dict):
                continue
            source_relative_path = file_entry.get("relative_path")
            if not isinstance(source_relative_path, str) or not source_relative_path.strip():
                continue

            raw_path = _resolve_manifest_entry_path(manifest_path, source_relative_path.strip())
            selected_path = _find_redacted_file(raw_path, backend_root=backend_root)
            if not selected_path.exists() or not selected_path.is_file():
                raise FileNotFoundError(
                    f"Dataset `{dataset_id}` source file missing: {source_relative_path}"
                )
            _ensure_supported(selected_path)

            extracted = extract_units(selected_path)
            spans = build_spans(
                extracted.units,
                max_chars=max_chars,
                min_chars=min_chars,
                overlap_chars=overlap_chars,
            )

            file_type = file_entry.get("file_type")
            if not isinstance(file_type, str) or not file_type.strip():
                file_type = selected_path.suffix.lower().lstrip(".")

            sha256 = _sha256_file(selected_path)
            document_id = _build_document_id(
                dataset_id=dataset_id,
                source_relative_path=source_relative_path,
                sha256=sha256,
            )
            file_char_cursor = 0

            for span in spans:
                chunk_start_char = file_char_cursor
                chunk_end_char = chunk_start_char + len(span.text)
                file_char_cursor = chunk_end_char

                locator = _build_locator(
                    locator_key=extracted.locator_key,
                    span=span,
                    source_relative_path=source_relative_path,
                    document_id=document_id,
                    chunk_start_char=chunk_start_char,
                    chunk_end_char=chunk_end_char,
                )

                meta: dict[str, Any] = {
                    "source_relative_path": source_relative_path,
                    "file_type": file_type,
                    "sha256": sha256,
                    "dataset_id": dataset_id,
                    "created_at": created_at,
                    "char_count": len(span.text),
                    "locator": locator,
                }
                if extracted.locator_key == "page":
                    meta["page_start"] = span.start
                    meta["page_end"] = span.end
                elif extracted.locator_key == "paragraph":
                    meta["paragraph_start"] = span.start
                    meta["paragraph_end"] = span.end
                elif extracted.locator_key == "slide":
                    meta["slide_start"] = span.start
                    meta["slide_end"] = span.end

                if span.title:
                    meta["title"] = span.title
                if span.heading:
                    meta["heading"] = span.heading

                records.append(
                    {
                        "chunk_id": f"ch-{dataset_id}-{sequence:06d}",
                        "text": span.text,
                        "meta": meta,
                    }
                )
                sequence += 1

        dataset_chunks[dataset_id] = records

        lengths = [len(record["text"]) for record in records]
        total_chars = sum(lengths)
        chunk_count = len(records)
        per_dataset_stats[dataset_id] = {
            "chunk_count": chunk_count,
            "total_chars": total_chars,
            "avg_chars": round(total_chars / chunk_count, 2) if chunk_count else 0,
            "min_chars": min(lengths) if lengths else 0,
            "max_chars": max(lengths) if lengths else 0,
        }

    stats = {
        "version": "v1",
        "created_at": created_at,
        "max_chars": max_chars,
        "min_chars": min_chars,
        "overlap_chars": overlap_chars,
        "datasets": per_dataset_stats,
        "totals": {
            "dataset_count": len(dataset_chunks),
            "chunk_count": sum(entry["chunk_count"] for entry in per_dataset_stats.values()),
            "total_chars": sum(entry["total_chars"] for entry in per_dataset_stats.values()),
        },
    }
    return dataset_chunks, stats


def write_outputs(
    *,
    out_dir: Path,
    dataset_chunks: dict[str, list[dict[str, Any]]],
    stats: dict[str, Any],
) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    for dataset_id, records in dataset_chunks.items():
        out_path = out_dir / f"{dataset_id}.chunks.v1.jsonl"
        with out_path.open("w", encoding="utf-8") as file_obj:
            for record in records:
                file_obj.write(json.dumps(record, ensure_ascii=False))
                file_obj.write("\n")

    stats_path = out_dir / "chunk_stats.v1.json"
    stats_path.write_text(
        json.dumps(stats, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build T3 chunks from dataset manifest files.")
    parser.add_argument(
        "--manifest",
        type=Path,
        default=DEFAULT_MANIFEST_PATH,
        help=f"Path to manifest.v1.yaml (default: {DEFAULT_MANIFEST_PATH})",
    )
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=DEFAULT_OUT_DIR,
        help=f"Output directory for chunk JSONL and stats (default: {DEFAULT_OUT_DIR})",
    )
    parser.add_argument(
        "--max-chars",
        type=int,
        default=1200,
        help="Maximum characters per chunk (default: 1200)",
    )
    parser.add_argument(
        "--min-chars",
        type=int,
        default=200,
        help="Target minimum characters per chunk when mergeable (default: 200)",
    )
    parser.add_argument(
        "--overlap-chars",
        type=int,
        default=0,
        help="Optional overlap characters between adjacent chunks (default: 0)",
    )
    parser.add_argument(
        "--dataset-id",
        type=str,
        default=None,
        help="Optional dataset_id filter; when set only one dataset is processed.",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Build in-memory stats only, skip writing JSONL/stats files.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        dataset_chunks, stats = build_dataset_chunks(
            manifest_path=args.manifest,
            max_chars=args.max_chars,
            min_chars=args.min_chars,
            overlap_chars=args.overlap_chars,
            dataset_id_filter=args.dataset_id,
        )
    except Exception as exc:
        print(f"build_chunks failed: {exc}", file=sys.stderr)
        return 1

    if args.dry_run:
        print(json.dumps(stats, ensure_ascii=False, indent=2))
    else:
        try:
            write_outputs(out_dir=args.out_dir, dataset_chunks=dataset_chunks, stats=stats)
        except Exception as exc:
            print(f"Failed to write chunk outputs: {exc}", file=sys.stderr)
            return 1

    total_chunks = stats["totals"]["chunk_count"]
    dataset_count = stats["totals"]["dataset_count"]
    out_dir = args.out_dir.resolve()
    suffix = " (dry-run, files not written)" if args.dry_run else ""
    print(f"Processed datasets={dataset_count}, total_chunks={total_chunks}, output={out_dir}{suffix}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
