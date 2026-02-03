import shutil
import os
import uuid
from pathlib import Path
from fastapi import UploadFile
import fitz 
from docx import Document
from pptx import Presentation

UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

def save_upload_file(upload_file: UploadFile) -> str:
    file_ext = os.path.splitext(upload_file.filename)[1].lower()
    unique_name = f"{uuid.uuid4()}{file_ext}"
    file_path = UPLOAD_DIR / unique_name
    
    with file_path.open("wb") as buffer:
        shutil.copyfileobj(upload_file.file, buffer)
        
    return str(file_path)

def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    if ext == ".pdf":
        with fitz.open(file_path) as doc:
            for page in doc: text += page.get_text()
    elif ext == ".docx":
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text += shape.text + "\n"
    return text.strip()